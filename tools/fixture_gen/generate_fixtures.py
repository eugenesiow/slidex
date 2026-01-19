from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[2]
FIXTURES = ROOT / "tests" / "fixtures" / "generated" / "simple"


def title_and_content() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = "Hello World"
    slide.placeholders[1].text = "This is text"

    prs.save(FIXTURES / "title_and_content.pptx")


def text_only() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    left = top = width = height = 1000000
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.text = "First line"
    p = tf.add_paragraph()
    p.text = "Second line"
    p = tf.add_paragraph()
    p.text = "Hello {{name}}"

    prs.save(FIXTURES / "text_only.pptx")


def blank_template() -> None:
    assets = ROOT / "assets"
    assets.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(assets / "blank.pptx")


def main() -> None:
    title_and_content()
    text_only()
    blank_template()


if __name__ == "__main__":
    main()
